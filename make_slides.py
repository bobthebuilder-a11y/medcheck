from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
BLUE = RGBColor(0x26, 0x63, 0xEB)       # blue-600
DARK = RGBColor(0x11, 0x18, 0x27)       # near black
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF9, 0xFA, 0xFB)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x77, 0x06)
PURPLE = RGBColor(0x79, 0x16, 0xF6)
LIGHT_BLUE = RGBColor(0xEF, 0xF6, 0xFF)

blank_layout = prs.slide_layouts[6]  # blank

def add_rect(slide, l, t, w, h, fill=None, line=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line and line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def new_slide():
    return prs.slides.add_slide(blank_layout)

# ─────────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────────
slide = new_slide()
# Full background
add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
# Blue accent bar left
add_rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)
# Blue accent circle (decorative)
add_rect(slide, 9.5, 1.5, 4.5, 4.5, fill=RGBColor(0x1e, 0x40, 0xaf))

# Icon
add_text(slide, "🔬", 1.0, 1.2, 2, 1.5, size=72, align=PP_ALIGN.LEFT)
# Title
add_text(slide, "MedCheck", 1.0, 2.8, 8, 1.4, size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
# Subtitle
add_text(slide, "AI-Powered Health Misinformation Detection", 1.0, 4.1, 9, 0.6, size=24, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.LEFT)
# Line
add_rect(slide, 1.0, 4.85, 5, 0.04, fill=BLUE)
# Details
add_text(slide, "David Xiao  ·  ACP Student AI Championship 2026", 1.0, 5.1, 9, 0.45, size=16, color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.LEFT)
add_text(slide, "SDG 3: Good Health & Well-Being  ·  SDG 16: Peace, Justice & Strong Institutions", 1.0, 5.55, 10, 0.45, size=14, color=RGBColor(0x6B, 0x72, 0x80), align=PP_ALIGN.LEFT)
add_text(slide, "medcheck-murex.vercel.app", 1.0, 6.1, 6, 0.4, size=14, color=BLUE, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────
# SLIDE 2: The Problem
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill=DARK)
add_rect(slide, 0, 1.08, 13.33, 0.04, fill=BLUE)
add_text(slide, "Health Misinformation Is Lethal", 0.5, 0.18, 12, 0.75, size=32, bold=True, color=WHITE)

# 4 stat boxes
stats = [
    ("6×", "Faster than corrections,\nfalse claims spread online"),
    ("30M+", "Americans making health decisions\nfrom unverified social media"),
    ("1 in 3", "Americans have acted on\nunverified health advice online"),
    ("3 sec\nvs 20 min", "To read a false claim vs.\nto manually fact-check it"),
]
cols = [0.4, 3.55, 6.7, 9.85]
for i, (val, label) in enumerate(stats):
    x = cols[i]
    box = add_rect(slide, x, 1.3, 3.0, 2.8, fill=WHITE)
    box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    box.line.width = Pt(1)
    add_text(slide, val, x+0.1, 1.45, 2.8, 1.1, size=36, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x+0.1, 2.55, 2.8, 0.9, size=13, color=GRAY, align=PP_ALIGN.CENTER)

# Quote
add_rect(slide, 0.4, 4.35, 12.5, 0.9, fill=RGBColor(0xEF, 0xF6, 0xFF))
add_text(slide, '"By the time a correction reaches someone, they\'ve already shared the false claim."', 0.6, 4.45, 12.1, 0.7, size=16, color=BLUE, align=PP_ALIGN.CENTER)

# Bottom note
add_text(slide, "Sources: WHO Infodemic Report · Johns Hopkins Bloomberg School of Public Health · Pew Research Center", 0.4, 5.5, 12.5, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Why existing fail
add_text(slide, "❌  Snopes / PolitiFact — you must already be skeptical to search", 0.5, 6.05, 12, 0.4, size=13, color=DARK)
add_text(slide, "❌  Social media flags — applied after viral spread, never explain why", 0.5, 6.48, 12, 0.4, size=13, color=DARK)

# ─────────────────────────────────────────────
# SLIDE 3: Solution Overview
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill=BLUE)
add_text(slide, "Introducing MedCheck", 0.5, 0.18, 12, 0.75, size=32, bold=True, color=WHITE)

add_text(slide, "Paste any health claim. Get a structured, sourced, calibrated verdict in seconds.", 0.5, 1.25, 12.3, 0.5, size=17, color=GRAY, align=PP_ALIGN.CENTER)

features = [
    ("✅ / ❌ / ⚠️ / ❓", "4-Verdict System", "TRUE · FALSE · MISLEADING · UNVERIFIABLE\nNot just binary — reflects real-world nuance"),
    ("📊", "Confidence Score", "0–100% calibrated uncertainty\nThe AI says 'I'm not sure' when it isn't"),
    ("🔍", "Claim Breakdown", "Each assertion evaluated separately\nBreaks compound claims apart"),
    ("📚", "Real Citations", "CDC · WHO · NIH · PubMed\nSo you can verify independently"),
    ("⚡", "Political Charge Flag", "Flags where science meets politics\nConnected to ongoing research"),
    ("🕐", "History & Sharing", "Persists last 20 checks\nShare results with one tap"),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.28
    y = 1.95 + row * 2.45
    box = add_rect(slide, x, y, 4.0, 2.2, fill=WHITE)
    box.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
    box.line.width = Pt(1)
    add_text(slide, icon, x+0.15, y+0.12, 0.8, 0.55, size=22, align=PP_ALIGN.LEFT)
    add_text(slide, title, x+0.15, y+0.65, 3.7, 0.4, size=14, bold=True, color=DARK)
    add_text(slide, desc, x+0.15, y+1.05, 3.7, 1.0, size=11, color=GRAY)

# ─────────────────────────────────────────────
# SLIDE 4: Architecture
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x1e, 0x3a, 0x8a))
add_text(slide, "🏗️  How It Works — Architecture", 0.5, 0.18, 12, 0.75, size=32, bold=True, color=WHITE)

# Pipeline
steps = ["User inputs claim", "AI decomposes\nassertions", "Evidence synthesis\nCDC · WHO · PubMed", "Confidence\ncalibration", "Structured\nJSON output", "React frontend\nrenders result"]
arrows = ["→", "→", "→", "→", "→"]
box_w = 1.85
gap = 0.18
start_x = 0.3

for i, step in enumerate(steps):
    x = start_x + i * (box_w + gap + 0.25)
    box = add_rect(slide, x, 1.3, box_w, 1.5, fill=RGBColor(0x1e, 0x40, 0xaf))
    add_text(slide, step, x+0.08, 1.38, box_w-0.16, 1.35, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        add_text(slide, "→", x+box_w+0.03, 1.75, 0.25, 0.5, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Tech stack
add_rect(slide, 0.3, 3.1, 12.7, 0.04, fill=RGBColor(0x37, 0x41, 0x51))
add_text(slide, "Tech Stack", 0.3, 3.25, 3, 0.4, size=14, bold=True, color=RGBColor(0x9C, 0xA3, 0xAF))
tech = [
    ("⚛️  React + TypeScript", "Type-safe modern frontend"),
    ("🤖  Llama 4 via Groq", "Fast LLM inference, free tier"),
    ("🎨  Tailwind CSS", "Utility-first responsive styling"),
    ("▲  Vercel", "Global CDN deployment"),
    ("📦  Vite", "Lightning-fast build tool"),
]
for i, (name, desc) in enumerate(tech):
    x = 0.3 + i * 2.6
    add_text(slide, name, x, 3.75, 2.5, 0.4, size=12, bold=True, color=WHITE)
    add_text(slide, desc, x, 4.15, 2.5, 0.4, size=11, color=GRAY)

# Key innovation box
add_rect(slide, 0.3, 4.75, 12.7, 1.65, fill=RGBColor(0x0c, 0x24, 0x61))
add_text(slide, "🔑  Key Innovation: Structured Prompting for Calibrated Fact-Checking", 0.6, 4.9, 12.1, 0.45, size=15, bold=True, color=WHITE)
add_text(slide, 'The system prompt explicitly instructs the AI: "Never express false confidence. A low-confidence honest answer is better than a high-confidence wrong answer." This is not just calling an API — it\'s engineered prompting for a specific analytical task with defined output structure.', 0.6, 5.4, 12.1, 0.9, size=12, color=RGBColor(0x93, 0xC5, 0xFD))

# ─────────────────────────────────────────────
# SLIDE 5: The 4 Verdicts
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill=DARK)
add_text(slide, "Honest Verdicts for a Complex World", 0.5, 0.18, 12, 0.75, size=32, bold=True, color=WHITE)

verdicts = [
    ("✅", "TRUE", "Scientific evidence clearly supports this claim", GREEN, RGBColor(0xDC, 0xFC, 0xE7)),
    ("❌", "FALSE", "Scientific evidence clearly contradicts this claim", RED, RGBColor(0xFE, 0xE2, 0xE2)),
    ("⚠️", "MISLEADING", "Contains partial truths but creates a false overall impression", AMBER, RGBColor(0xFE, 0xF3, 0xC7)),
    ("❓", "UNVERIFIABLE", "Insufficient scientific consensus to evaluate definitively", GRAY, RGBColor(0xF3, 0xF4, 0xF6)),
]

for i, (icon, label, desc, color, bg) in enumerate(verdicts):
    y = 1.25 + i * 1.42
    box = add_rect(slide, 0.4, y, 12.5, 1.3, fill=bg)
    box.line.color.rgb = color
    box.line.width = Pt(2)
    add_text(slide, icon, 0.55, y+0.1, 0.7, 0.7, size=28, align=PP_ALIGN.LEFT)
    add_text(slide, label, 1.35, y+0.1, 2.5, 0.5, size=20, bold=True, color=color)
    add_text(slide, desc, 1.35, y+0.6, 8.5, 0.5, size=13, color=DARK)

# Callout for MISLEADING
add_rect(slide, 0.4, 6.85, 12.5, 0.55, fill=RGBColor(0xFF, 0xF7, 0xE0))
add_text(slide, '💡  MISLEADING is the most important verdict — most dangerous health misinfo isn\'t false, it\'s selectively true. "Natural immunity is always better than vaccines" — partially true, dangerously misleading overall.', 0.6, 6.88, 12.1, 0.48, size=12, color=AMBER)

# ─────────────────────────────────────────────
# SLIDE 6: Why AI Is the Right Tool
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(slide, 0, 0, 13.33, 1.1, fill=BLUE)
add_text(slide, "🤖  Why AI Is the Right Tool for This Problem", 0.5, 0.18, 12, 0.75, size=30, bold=True, color=WHITE)

add_text(slide, "AI is not the point. Impact is. AI is the only tool fast enough and scalable enough to make this work.", 0.5, 1.2, 12.3, 0.5, size=15, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.CENTER)

reasons = [
    ("🌍", "Scale", "Millions of new health claims are posted daily. Human fact-checkers can't keep up. AI never sleeps and scales to any volume."),
    ("🧩", "Decomposition", "AI can break compound claims into individual assertions automatically — \"vaccines cause autism AND damage immunity\" is two claims, not one."),
    ("⚖️", "Calibrated Uncertainty", "AI can express degrees of confidence — not just true/false. A tool that says \"I'm not sure\" is more trustworthy than one that's always confident."),
    ("📚", "Evidence Synthesis", "Cross-references CDC, WHO, NIH, and PubMed simultaneously in seconds — a task that takes a human 20+ minutes."),
    ("🔓", "Accessibility", "No medical background required. Anyone can use it. Free. Works in any browser. Zero installation."),
]

for i, (icon, title, desc) in enumerate(reasons):
    col = i % 3 if i < 3 else (i - 3)
    row = 0 if i < 3 else 1
    x = 0.3 + (col if i < 3 else col + 0.65) * 4.2
    if i >= 3:
        x = 0.3 + (i - 3) * 6.35
    y = 1.9 + row * 2.45
    w = 4.0 if i < 3 else 6.0
    box = add_rect(slide, x, y, w, 2.2, fill=RGBColor(0x1f, 0x2a, 0x3d))
    add_text(slide, icon, x+0.15, y+0.12, 0.8, 0.55, size=22)
    add_text(slide, title, x+0.15, y+0.65, w-0.3, 0.4, size=14, bold=True, color=WHITE)
    add_text(slide, desc, x+0.15, y+1.1, w-0.3, 1.0, size=11, color=GRAY)

add_rect(slide, 0.3, 6.8, 12.7, 0.55, fill=RGBColor(0x1e, 0x3a, 0x8a))
add_text(slide, '💡  Key insight: The cognitive task of fact-checking — finding evidence, weighing it, synthesizing a verdict — is exactly what LLMs are trained to do. MedCheck harnesses that for public health.', 0.5, 6.85, 12.3, 0.45, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 7: Impact & SDGs
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill=GREEN)
add_text(slide, "🌍  Real-World Impact & SDG Alignment", 0.5, 0.18, 12, 0.75, size=32, bold=True, color=WHITE)

# SDG 3 box
add_rect(slide, 0.4, 1.25, 6.0, 2.7, fill=WHITE)
add_rect(slide, 0.4, 1.25, 6.0, 0.5, fill=GREEN)
add_text(slide, "SDG 3 — Good Health & Well-Being", 0.55, 1.3, 5.7, 0.4, size=14, bold=True, color=WHITE)
sdg3 = [
    "Reduces health decisions made on false information",
    "Disproportionate impact on underserved communities with limited healthcare access",
    "Scalable: free, runs in any browser, zero installation",
    "Every accurate claim checked = potential harm prevented",
]
for i, s in enumerate(sdg3):
    add_text(slide, f"• {s}", 0.55, 1.9+i*0.5, 5.7, 0.45, size=12, color=DARK)

# SDG 16 box
add_rect(slide, 6.9, 1.25, 6.0, 2.7, fill=WHITE)
add_rect(slide, 6.9, 1.25, 6.0, 0.5, fill=BLUE)
add_text(slide, "SDG 16 — Peace, Justice & Strong Institutions", 7.05, 1.3, 5.7, 0.4, size=14, bold=True, color=WHITE)
sdg16 = [
    "Informed public = functional health policy that works",
    "Counters coordinated health misinformation campaigns",
    "Builds media literacy, not AI dependency",
    "Open source = accessible to all, not just tech-savvy users",
]
for i, s in enumerate(sdg16):
    add_text(slide, f"• {s}", 7.05, 1.9+i*0.5, 5.7, 0.45, size=12, color=DARK)

# Scale path
add_rect(slide, 0.4, 4.15, 12.5, 0.4, fill=WHITE)
add_text(slide, "Scale Path:", 0.6, 4.18, 1.5, 0.35, size=13, bold=True, color=DARK)
scale = ["V1 (Now): Web App", "→", "V2: Browser Extension\n(auto-flags claims)", "→", "V3: Open-Source API\nfor newsrooms & health orgs", "→", "MediaLens Foundation\n(nonprofit)"]
positions = [2.2, 4.4, 4.6, 7.2, 7.4, 10.0, 10.2]
for i, (text, pos) in enumerate(zip(scale, positions)):
    is_arrow = text == "→"
    add_text(slide, text, pos, 4.18, 1.7 if not is_arrow else 0.3, 0.35, size=11 if not is_arrow else 16, color=BLUE if is_arrow else DARK, bold=is_arrow, align=PP_ALIGN.CENTER if is_arrow else PP_ALIGN.LEFT)

# Who uses it
add_rect(slide, 0.4, 4.75, 12.5, 2.55, fill=WHITE)
add_text(slide, "Who Uses MedCheck", 0.6, 4.85, 5, 0.4, size=14, bold=True, color=DARK)
users = [
    ("🙋", "General Public", "Anyone reading health content on social media"),
    ("👨‍👩‍👧", "Caregivers & Parents", "Making health decisions for families"),
    ("🎓", "Students", "Learning to evaluate health information critically"),
    ("🏥", "Public Health Workers", "Quickly screening claims for outreach programs"),
]
for i, (icon, role, desc) in enumerate(users):
    x = 0.4 + i * 3.22
    add_text(slide, icon, x+0.1, 5.35, 0.5, 0.5, size=20)
    add_text(slide, role, x+0.1, 5.85, 3.0, 0.35, size=12, bold=True, color=DARK)
    add_text(slide, desc, x+0.1, 6.2, 3.0, 0.5, size=11, color=GRAY)

# ─────────────────────────────────────────────
# SLIDE 8: Ethical Considerations
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, fill=AMBER)
add_text(slide, "⚠️  Ethical Considerations — We Thought About the Risks", 0.5, 0.18, 12, 0.75, size=28, bold=True, color=WHITE)

risks = [
    ("AI Overconfidence", "AI might sound confident when wrong — dangerous in health contexts", "Confidence scores shown explicitly. Low confidence verdicts are a FEATURE, not a bug."),
    ("User Dependency", "People might blindly trust AI instead of thinking critically", "Tool explains WHY, not just WHAT. Builds literacy, not reliance. Citations always provided."),
    ("False Authority", "Users might treat AI as a medical oracle", 'Explicit medical disclaimer on EVERY result: "Not a substitute for medical advice."'),
    ("Bias on Charged Topics", "AI may be less reliable on politically contested claims", "Political Charge Flag (⚡) warns users explicitly. Connected to our research paper."),
    ("Hallucination Risk", "AI might cite sources that don't exist or say wrong things", "Citations linked so users can verify. Confidence score reflects uncertainty. Encourage verification."),
]

for i, (risk, problem, mitigation) in enumerate(risks):
    y = 1.25 + i * 1.2
    # Risk label
    add_rect(slide, 0.4, y, 2.8, 1.05, fill=RGBColor(0xFF, 0xF3, 0xCD))
    add_text(slide, risk, 0.55, y+0.1, 2.5, 0.4, size=13, bold=True, color=AMBER)
    add_text(slide, problem, 0.55, y+0.5, 2.5, 0.5, size=10, color=DARK)
    # Arrow
    add_text(slide, "→", 3.25, y+0.3, 0.4, 0.5, size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    # Mitigation
    add_rect(slide, 3.7, y, 9.2, 1.05, fill=RGBColor(0xF0, 0xFD, 0xF4))
    add_text(slide, "✅  " + mitigation, 3.85, y+0.25, 8.9, 0.6, size=12, color=GREEN)

add_rect(slide, 0.4, 7.1, 12.5, 0.3, fill=RGBColor(0xFE, 0xE2, 0xE2))
add_text(slide, '"This tool can be wrong. The goal is to help people think more carefully — not to replace thinking."', 0.6, 7.12, 12.1, 0.26, size=11, color=RED, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 9: Future Directions
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(slide, 0, 0, 13.33, 1.1, fill=BLUE)
add_text(slide, "🚀  Future Directions", 0.5, 0.18, 12, 0.75, size=32, bold=True, color=WHITE)

timeline = [
    ("Now", "V1: Web App", "Working prototype live\nmedcheck-murex.vercel.app\n4-verdict system + citations\nHistory & sharing built in", BLUE),
    ("1-3 mo", "V2: Browser Extension", "Claims flagged automatically\nas you browse social media\nNo copy-paste required\nMeets people where they are", RGBColor(0x06, 0x95, 0xc5)),
    ("6 mo", "Research Paper", "150 claims × 3 models\nAI bias measurement\nProfessor mentorship\nTarget: JEI publication", PURPLE),
    ("1 yr", "Multi-Domain + API", "Climate, politics, history\nOpen API for newsrooms\nPublic health partnerships\nSchool library deployment", GREEN),
    ("2+ yr", "MediaLens Foundation", "Open-source nonprofit\nFree for everyone\nMedia literacy curriculum\nPlatform integration", AMBER),
]

for i, (time, title, desc, color) in enumerate(timeline):
    x = 0.3 + i * 2.6
    # Timeline dot
    add_rect(slide, x+0.95, 1.2, 0.7, 0.7, fill=color)
    add_text(slide, time, x, 1.3, 2.5, 0.45, size=11, color=color, align=PP_ALIGN.CENTER, bold=True)
    # Connector line (except last)
    if i < len(timeline)-1:
        add_rect(slide, x+1.6, 1.5, 1.0, 0.06, fill=RGBColor(0x37, 0x41, 0x51))
    # Card
    add_rect(slide, x, 2.1, 2.5, 4.6, fill=RGBColor(0x1f, 0x2a, 0x3d))
    add_rect(slide, x, 2.1, 2.5, 0.5, fill=color)
    add_text(slide, title, x+0.1, 2.13, 2.3, 0.43, size=12, bold=True, color=WHITE)
    add_text(slide, desc, x+0.1, 2.7, 2.3, 3.9, size=11, color=RGBColor(0x9C, 0xA3, 0xAF))

# ─────────────────────────────────────────────
# SLIDE 10: Conclusion
# ─────────────────────────────────────────────
slide = new_slide()
add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
add_rect(slide, 0, 0, 0.08, 7.5, fill=BLUE)

add_text(slide, "🔬", 1.0, 0.5, 1.5, 1.2, size=64)
add_text(slide, "The Problem Won't Wait.", 1.0, 1.5, 10, 0.9, size=48, bold=True, color=WHITE)
add_text(slide, "Health misinformation kills. AI can help. We built it.", 1.0, 2.5, 10, 0.5, size=22, color=RGBColor(0x93, 0xC5, 0xFD))

add_rect(slide, 1.0, 3.2, 11.0, 0.04, fill=RGBColor(0x37, 0x41, 0x51))

checkmarks = [
    "Any health claim analyzed in seconds — true, false, misleading, or unverifiable",
    "Honest confidence calibration — the AI admits uncertainty rather than faking confidence",
    "Explains WHY, not just WHAT — building real media literacy, not AI dependency",
    "Real source citations — verifiable, not blind trust",
    "Political charge flag — connected to original research on AI bias",
    "Free. Open. Accessible to everyone with internet access.",
]
for i, c in enumerate(checkmarks):
    add_text(slide, "✅  " + c, 1.0, 3.45+i*0.52, 11.0, 0.48, size=13, color=WHITE)

add_rect(slide, 1.0, 6.7, 5.5, 0.65, fill=BLUE)
add_text(slide, "🌐  medcheck-murex.vercel.app", 1.15, 6.75, 5.2, 0.55, size=14, bold=True, color=WHITE)
add_rect(slide, 7.0, 6.7, 5.5, 0.65, fill=RGBColor(0x1e, 0x3a, 0x8a))
add_text(slide, "💻  github.com/bobthebuilder-a11y/medcheck", 7.15, 6.75, 5.2, 0.55, size=13, color=WHITE)

# Save
prs.save('/Users/xiaofamily/medcheck/MedCheck_ACP_2026.pptx')
print("✅ Saved: MedCheck_ACP_2026.pptx")
